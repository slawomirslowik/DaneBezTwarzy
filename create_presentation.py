"""
Skrypt do automatycznego tworzenia prezentacji PowerPoint dla projektu Dane Bez Twarzy
Wymaga: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    """Tworzy prezentację PowerPoint z 5 slajdami o projekcie Dane Bez Twarzy."""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Kolory
    PRIMARY_COLOR = RGBColor(102, 126, 234)  # #667eea
    SECONDARY_COLOR = RGBColor(118, 75, 162)  # #764ba2
    DARK_COLOR = RGBColor(51, 51, 51)
    LIGHT_BG = RGBColor(248, 249, 250)
    
    # ========== SLAJD 1: Tytuł ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Pusty layout
    
    # Gradient tło (aproksymacja przez prostokąt)
    background = slide1.shapes.add_shape(
        1,  # Rectangle
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY_COLOR
    background.line.fill.background()
    
    # Tytuł główny
    title_box = slide1.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Dane Bez Twarzy"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Podtytuł
    subtitle_box = slide1.shapes.add_textbox(
        Inches(1), Inches(3.7), Inches(8), Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Automatyczna Anonimizacja Danych Osobowych"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Informacje dodatkowe
    info_box = slide1.shapes.add_textbox(
        Inches(1), Inches(5), Inches(8), Inches(1.5)
    )
    info_frame = info_box.text_frame
    info_frame.text = "Python • RODO/GDPR • Open Source\nBiblioteka + CLI + Standalone EXE"
    for paragraph in info_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(230, 230, 230)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # ========== SLAJD 2: Problem i Rozwiązanie ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    
    title = slide2.shapes.title
    title.text = "Problem i Rozwiązanie"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    title.text_frame.paragraphs[0].font.bold = True
    
    content = slide2.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Problem
    p = tf.add_paragraph()
    p.text = "Problem"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_COLOR
    p.space_after = Pt(10)
    
    problems = [
        "Dane osobowe w dokumentach wymagają ochrony zgodnie z RODO/GDPR",
        "Ręczna anonimizacja setek plików jest czasochłonna i błędogenna",
        "Brak uniwersalnego narzędzia dla różnych formatów (DOCX, PDF, XLSX)"
    ]
    
    for problem in problems:
        p = tf.add_paragraph()
        p.text = problem
        p.level = 1
        p.font.size = Pt(18)
        p.space_after = Pt(8)
    
    # Rozwiązanie
    p = tf.add_paragraph()
    p.text = "Rozwiązanie"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.space_before = Pt(20)
    p.space_after = Pt(10)
    
    solutions = [
        "Automatyczne wykrywanie 25+ typów danych osobowych",
        "Wsparcie dla 5+ formatów plików (TXT, DOCX, PDF, XLSX, CSV)",
        "3 interfejsy: Biblioteka Python, CLI, Standalone EXE"
    ]
    
    for solution in solutions:
        p = tf.add_paragraph()
        p.text = solution
        p.level = 1
        p.font.size = Pt(18)
        p.space_after = Pt(8)
    
    # ========== SLAJD 3: Architektura i Technologie ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide3.shapes.title
    title.text = "Architektura i Technologie"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    title.text_frame.paragraphs[0].font.bold = True
    
    content = slide3.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Detektory
    p = tf.add_paragraph()
    p.text = "4 Detektory Działające Równolegle"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_COLOR
    p.space_after = Pt(10)
    
    detectors = [
        ("PlaceholderDetector", "Wykrywa [name], [email], [pesel] - pewność 1.0"),
        ("RegexDetector", "PESEL, NIP, email, telefon - pewność 0.8-0.95"),
        ("PolishDetector", "Polskie wzorce: adresy, kody pocztowe - pewność 0.7-0.9"),
        ("NLPDetector", "spaCy NER dla imion/nazwisk - pewność 0.6-0.95"),
        ("LLMDetector", "PLLUM (12B parametrów) - najwyższa dokładność")
    ]
    
    for name, desc in detectors:
        p = tf.add_paragraph()
        p.text = f"{name}: {desc}"
        p.level = 1
        p.font.size = Pt(16)
        p.space_after = Pt(6)
    
    # Technologie
    p = tf.add_paragraph()
    p.text = "Stack Technologiczny"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.space_before = Pt(20)
    p.space_after = Pt(10)
    
    tech = [
        "Python 3.12 • spaCy 3.7 • langchain-openai",
        "PyPDF2, python-docx, openpyxl (procesory formatów)",
        "matplotlib, plotly (wizualizacja raportów)",
        "Nuitka (kompilacja do standalone EXE)"
    ]
    
    for t in tech:
        p = tf.add_paragraph()
        p.text = t
        p.level = 1
        p.font.size = Pt(16)
        p.space_after = Pt(6)
    
    # ========== SLAJD 4: Funkcjonalności ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide4.shapes.title
    title.text = "Kluczowe Funkcjonalności"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    title.text_frame.paragraphs[0].font.bold = True
    
    content = slide4.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    features = [
        ("6 Metod Anonimizacji", [
            "Maskowanie (***), Pseudonimizacja (Osoba_A)",
            "Entity ([name] [surname]), Haszowanie, Szyfrowanie, Redakcja"
        ]),
        ("Przetwarzanie Wsadowe", [
            "Pojedyncze pliki lub całe katalogi (rekurencyjnie)",
            "Wzorce plików: *.txt, *.docx, *.xlsx, *.pdf"
        ]),
        ("Raportowanie", [
            "3 formaty: JSON (dane), HTML (interaktywne wykresy), PDF (druk)",
            "Statystyki: liczba encji, typy, pewność, czas wykonania",
            "Automatyczna rotacja logów (5 × 10 MB)"
        ]),
        ("Chunking dla LLM", [
            "Automatyczny podział dużych plików (3000 znaków/fragment)",
            "Nakładanie fragmentów (200 znaków) - eliminacja utraty kontekstu",
            "Deduplikacja wykrytych encji"
        ])
    ]
    
    for feature_title, feature_points in features:
        p = tf.add_paragraph()
        p.text = feature_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = DARK_COLOR
        p.space_after = Pt(6)
        
        for point in feature_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 1
            p.font.size = Pt(15)
            p.space_after = Pt(4)
    
    # ========== SLAJD 5: Demo i Podsumowanie ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide5.shapes.title
    title.text = "Demo i Wyniki"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    title.text_frame.paragraphs[0].font.bold = True
    
    content = slide5.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Przykład użycia
    p = tf.add_paragraph()
    p.text = "Przykład CLI"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_COLOR
    p.space_after = Pt(10)
    
    # Kod w ramce
    code_box = slide5.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(2.2), Inches(9), Inches(1.2)
    )
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(40, 44, 52)
    code_box.line.color.rgb = PRIMARY_COLOR
    code_box.line.width = Pt(2)
    
    code_text = code_box.text_frame
    code_text.text = "dane-bez-twarzy anonymize cv.docx -o cv_anon.docx \\\n  --method entity --use-nlp --use-llm \\\n  --llm-api-key \"klucz\" --add-report raport --report-format all -v"
    code_text.word_wrap = True
    for paragraph in code_text.paragraphs:
        paragraph.font.name = "Consolas"
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = RGBColor(171, 178, 191)
    
    # Wyniki
    p = tf.add_paragraph()
    p.text = "Osiągnięte Wyniki"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.space_before = Pt(100)
    p.space_after = Pt(10)
    
    results = [
        "✓ Wykryto 25+ typów danych osobowych (PESEL, NIP, email, telefon, imiona)",
        "✓ Chunking LLM obsługuje pliki >100 KB bez błędów 413/400",
        "✓ Raport HTML z wykresami Plotly (interaktywny) i PDF (printable)",
        "✓ Standalone EXE (~50 MB LITE, ~600 MB FULL z NLP)",
        "✓ Format ENTITY zgodny z nask_train (orig.txt)"
    ]
    
    for result in results:
        p = tf.add_paragraph()
        p.text = result
        p.level = 1
        p.font.size = Pt(16)
        p.space_after = Pt(6)
    
    # Footer na ostatnim slajdzie
    footer_box = slide5.shapes.add_textbox(
        Inches(0.5), Inches(6.8), Inches(9), Inches(0.5)
    )
    footer_frame = footer_box.text_frame
    footer_frame.text = "GitHub: slawomirslowik/DaneBezTwarzy • Email: semantis@int.pl • Licencja: MIT"
    p = footer_frame.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER
    
    # Zapisz prezentację
    output_path = "DaneBezTwarzy_Prezentacja.pptx"
    prs.save(output_path)
    print(f"✓ Prezentacja zapisana: {output_path}")
    print(f"  - Liczba slajdów: {len(prs.slides)}")
    print(f"  - Format: {prs.slide_width.inches:.1f}\" × {prs.slide_height.inches:.1f}\"")
    return output_path


if __name__ == "__main__":
    print("========================================")
    print("  Generator Prezentacji PowerPoint")
    print("  Dane Bez Twarzy")
    print("========================================")
    print()
    
    try:
        # Sprawdź czy python-pptx jest zainstalowane
        try:
            import pptx
            print("✓ Biblioteka python-pptx znaleziona")
        except ImportError:
            print("✗ Biblioteka python-pptx nie jest zainstalowana!")
            print()
            print("Instaluję python-pptx...")
            import subprocess
            subprocess.check_call(["pip", "install", "python-pptx"])
            print("✓ Zainstalowano python-pptx")
            print()
        
        # Twórz prezentację
        print("Tworzę prezentację...")
        output = create_presentation()
        
        print()
        print("========================================")
        print("  Sukces!")
        print("========================================")
        print()
        print(f"Otwórz plik: {output}")
        
    except Exception as e:
        print(f"✗ Błąd: {e}")
        import traceback
        traceback.print_exc()
